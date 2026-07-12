#!/usr/bin/env python3
import sys

print("Checking dependencies...")
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("\n[ERROR] The 'python-pptx' library is not installed.")
    print("Please install it by running the following command in your terminal:")
    print("    pip install python-pptx")
    sys.exit(1)

# Initialize presentation
prs = Presentation()
# Set widescreen format (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme Constants (Matching the HTML template)
COLOR_BG = RGBColor(13, 14, 21)         # Dark Purple-Grey (#0D0E15)
COLOR_PRIMARY = RGBColor(99, 102, 241)   # Purple-Blue (#6366F1)
COLOR_SECONDARY = RGBColor(168, 85, 247) # Violet (#A855F7)
COLOR_TEXT_MAIN = RGBColor(243, 244, 246) # White-Grey (#F3F4F6)
COLOR_TEXT_MUTED = RGBColor(156, 163, 175) # Muted Grey (#9CA3AF)
COLOR_CARD_BG = RGBColor(25, 26, 38)     # Dark card background (#191A26)

# Font configurations
FONT_TITLE = "Outfit"
FONT_BODY = "Inter"

def apply_slide_bg(slide):
    """Sets a solid dark background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="Online Saloon Management System (MSMS)"):
    """Adds a standard header section to a slide."""
    apply_slide_bg(slide)
    
    # Add top-right category category_text
    cat_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.4), Inches(4.4), Inches(0.4))
    cat_tf = cat_box.text_frame
    cat_tf.word_wrap = True
    p_cat = cat_tf.paragraphs[0]
    p_cat.alignment = PP_ALIGN.RIGHT
    run_cat = p_cat.add_run()
    run_cat.text = category_text
    run_cat.font.name = FONT_BODY
    run_cat.font.size = Pt(10)
    run_cat.font.bold = True
    run_cat.font.color.rgb = COLOR_PRIMARY
    
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(7.5), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p_title = title_tf.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = title_text
    run_title.font.name = FONT_TITLE
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = COLOR_TEXT_MAIN

def create_card_shape(slide, left, top, width, height):
    """Creates a stylized glassmorphic-like card base shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD_BG
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(1)
    return shape

# ==============================================================================
# SLIDE 1: Title Cover Slide
# ==============================================================================
slide_layout = prs.slide_layouts[6] # blank layout
slide1 = prs.slides.add_slide(slide_layout)
apply_slide_bg(slide1)

# Large Title
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.0))
tf1 = title_box.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
run1 = p1.add_run()
run1.text = "ONLINE SALOON MANAGEMENT SYSTEM"
run1.font.name = FONT_TITLE
run1.font.size = Pt(44)
run1.font.bold = True
run1.font.color.rgb = COLOR_TEXT_MAIN

# Subtitle
p1_sub = tf1.add_paragraph()
p1_sub.alignment = PP_ALIGN.CENTER
run1_sub = p1_sub.add_run()
run1_sub.text = "A PHP & MySQL Web Solution for Saloon Booking, Invoicing, and Staff Scheduling"
run1_sub.font.name = FONT_BODY
run1_sub.font.size = Pt(16)
run1_sub.font.color.rgb = COLOR_TEXT_MUTED

# Meta Card box
card1 = create_card_shape(slide1, Inches(4.166), Inches(4.5), Inches(5.0), Inches(1.8))
card_tf = card1.text_frame
card_tf.word_wrap = True
card_tf.margin_left = Inches(0.2)
card_tf.margin_right = Inches(0.2)
card_tf.margin_top = Inches(0.2)

p_meta = card_tf.paragraphs[0]
p_meta.alignment = PP_ALIGN.CENTER
run_meta = p_meta.add_run()
run_meta.text = "Project Viva Presentation\n"
run_meta.font.bold = True
run_meta.font.name = FONT_BODY
run_meta.font.size = Pt(13)
run_meta.font.color.rgb = COLOR_PRIMARY

p_meta2 = card_tf.add_paragraph()
p_meta2.alignment = PP_ALIGN.CENTER
run_meta2 = p_meta2.add_run()
run_meta2.text = "Full-Stack Web Engineering Report\n"
run_meta2.font.name = FONT_BODY
run_meta2.font.size = Pt(11)
run_meta2.font.color.rgb = COLOR_TEXT_MAIN

p_meta3 = card_tf.add_paragraph()
p_meta3.alignment = PP_ALIGN.CENTER
run_meta3 = p_meta3.add_run()
run_meta3.text = "Role Modules: Client Panel, Stylist Dashboard & Admin System"
run_meta3.font.name = FONT_BODY
run_meta3.font.size = Pt(10)
run_meta3.font.color.rgb = COLOR_TEXT_MUTED

# ==============================================================================
# SLIDE 2: Problem Statement & Objectives
# ==============================================================================
slide2 = prs.slides.add_slide(slide_layout)
add_header(slide2, "Problem Statement & Project Objectives")

# Left Column (Bullet List)
left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

def add_bullet(tf, bold_prefix, text_content):
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_b = p.add_run()
    run_b.text = "•  " + bold_prefix + "  "
    run_b.font.name = FONT_BODY
    run_b.font.size = Pt(14)
    run_b.font.bold = True
    run_b.font.color.rgb = COLOR_PRIMARY
    
    run_c = p.add_run()
    run_c.text = text_content
    run_c.font.name = FONT_BODY
    run_c.font.size = Pt(13)
    run_c.font.color.rgb = COLOR_TEXT_MAIN

add_bullet(ltf, "Manual Queue Overhead:", "Traditional booking causes long waiting list jams and bad scheduling.")
add_bullet(ltf, "Staff Schedule Inefficiencies:", "Stylists are booked without factoring specialty, causing workload imbalance.")
add_bullet(ltf, "Friction in Billing:", "Calculating invoices by hand is slow, error-prone, and hard to log.")
add_bullet(ltf, "Absence of Reports:", "Saloon managers lack records of sales and service parameters.")

# Right Column (Structured Objective Cards)
create_card_shape(slide2, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.2))
card2_tf = slide2.shapes[-1].text_frame
card2_tf.word_wrap = True
card2_tf.margin_top = Inches(0.15)
card2_tf.margin_left = Inches(0.2)
p_c1 = card2_tf.paragraphs[0]
r_c1_t = p_c1.add_run()
r_c1_t.text = "🎯  System Objectives\n\n"
r_c1_t.font.bold = True
r_c1_t.font.size = Pt(14)
r_c1_t.font.color.rgb = COLOR_SECONDARY
r_c1_d = p_c1.add_run()
r_c1_d.text = "Provide clients dynamic bookings, display only competent stylists specializing in customer selection, and compile invoices cleanly."
r_c1_d.font.size = Pt(11)
r_c1_d.font.color.rgb = COLOR_TEXT_MAIN

create_card_shape(slide2, Inches(7.0), Inches(4.3), Inches(5.5), Inches(2.2))
card3_tf = slide2.shapes[-1].text_frame
card3_tf.word_wrap = True
card3_tf.margin_top = Inches(0.15)
card3_tf.margin_left = Inches(0.2)
p_c2 = card3_tf.paragraphs[0]
r_c2_t = p_c2.add_run()
r_c2_t.text = "🚀  Impact Goals\n\n"
r_c2_t.font.bold = True
r_c2_t.font.size = Pt(14)
r_c2_t.font.color.rgb = COLOR_SECONDARY
r_c2_d = p_c2.add_run()
r_c2_d.text = "Reduces scheduling friction, coordinates styling staff schedules via unified dashboards, and delivers 100% paperless record management."
r_c2_d.font.size = Pt(11)
r_c2_d.font.color.rgb = COLOR_TEXT_MAIN

# ==============================================================================
# SLIDE 3: Tech Stack & Architecture
# ==============================================================================
slide3 = prs.slides.add_slide(slide_layout)
add_header(slide3, "Technology Stack & Core Architecture")

# Text Subtitle
sub_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.6))
sub_tf = sub_box.text_frame
sub_tf.word_wrap = True
sub_p = sub_tf.paragraphs[0]
sub_r = sub_p.add_run()
sub_r.text = "Built on the robust, lightweight, and open-source LAMP/WAMP client-server ecosystem."
sub_r.font.name = FONT_BODY
sub_r.font.size = Pt(14)
sub_r.font.color.rgb = COLOR_TEXT_MUTED

# 3 Columns for UI, Backend, DB
col_w = Inches(3.7)
col_h = Inches(4.5)
gaps = Inches(0.3)
start_x = Inches(0.8)
start_y = Inches(2.0)

# Card 1: Frontend
create_card_shape(slide3, start_x, start_y, col_w, col_h)
tf_c1 = slide3.shapes[-1].text_frame
tf_c1.word_wrap = True
tf_c1.margin_left = Inches(0.25)
tf_c1.margin_right = Inches(0.25)
tf_c1.margin_top = Inches(0.3)
p_card1 = tf_c1.paragraphs[0]
p_card1.add_run("💻  Frontend User Interface\n\n").font.size = Pt(16)
p_card1.paragraphs[0].runs[0].font.bold = True
p_card1.paragraphs[0].runs[0].font.color.rgb = COLOR_PRIMARY
p_card1.add_run("• HTML5 & CSS3: ").font.bold = True
p_card1.add_run("Semantic tags and structural styling grids.\n\n")
p_card1.add_run("• Bootstrap Framework: ").font.bold = True
p_card1.add_run("Guarantees clean, mobile-responsive grids and layout wrappers.\n\n")
p_card1.add_run("• JS & jQuery: ").font.bold = True
p_card1.add_run("Dynamic element controls, AJAX stylists filtering on booking.")
for run in p_card1.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# Card 2: Backend
create_card_shape(slide3, start_x + col_w + gaps, start_y, col_w, col_h)
tf_c2 = slide3.shapes[-1].text_frame
tf_c2.word_wrap = True
tf_c2.margin_left = Inches(0.25)
tf_c2.margin_right = Inches(0.25)
tf_c2.margin_top = Inches(0.3)
p_card2 = tf_c2.paragraphs[0]
p_card2.add_run("⚙️  Backend Processing\n\n").font.size = Pt(16)
p_card2.paragraphs[0].runs[0].font.bold = True
p_card2.paragraphs[0].runs[0].font.color.rgb = COLOR_PRIMARY
p_card2.add_run("• PHP Scripts: ").font.bold = True
p_card2.add_run("Handles server-side calculations, form validations, database queries, and role filters.\n\n")
p_card2.add_run("• Role Access Filters: ").font.bold = True
p_card2.add_run("Session verification blocks clients, stylists, and admins from accessing unauthorized portals.\n\n")
p_card2.add_run("• Security: ").font.bold = True
p_card2.add_run("Uses MD5 encryption for database passwords.")
for run in p_card2.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# Card 3: Database
create_card_shape(slide3, start_x + (col_w + gaps)*2, start_y, col_w, col_h)
tf_c3 = slide3.shapes[-1].text_frame
tf_c3.word_wrap = True
tf_c3.margin_left = Inches(0.25)
tf_c3.margin_right = Inches(0.25)
tf_c3.margin_top = Inches(0.3)
p_card3 = tf_c3.paragraphs[0]
p_card3.add_run("🗄️  MySQL Database Engine\n\n").font.size = Pt(16)
p_card3.paragraphs[0].runs[0].font.bold = True
p_card3.paragraphs[0].runs[0].font.color.rgb = COLOR_PRIMARY
p_card3.add_run("• Relational DB Schema: ").font.bold = True
p_card3.add_run("Normalizes data structures into specialized tables.\n\n")
p_card3.add_run("• Relational Integrity: ").font.bold = True
p_card3.add_run("Maps customer invoice IDs and matches specialty connections cleanly.\n\n")
p_card3.add_run("• SQL Procedures: ").font.bold = True
p_card3.add_run("Runs clean schema upgrades for customized modules (Client login & Stylist tracking tools).")
for run in p_card3.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# ==============================================================================
# SLIDE 4: Database Design
# ==============================================================================
slide4 = prs.slides.add_slide(slide_layout)
add_header(slide4, "Database Schema & Entity Models")

# Left text
left_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
tf4 = left_box4.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
p4.add_run("Relational Data Design (MySQL)\n\n").font.size = Pt(18)
p4.runs[0].font.bold = True
p4.runs[0].font.color.rgb = COLOR_PRIMARY

p4_b1 = tf4.add_paragraph()
p4_b1.space_after = Pt(10)
p4_b1.add_run("• 9 Total Tables: ").font.bold = True
p4_b1.add_run("Separates core system logic from metadata tables.\n\n")
p4_b2 = tf4.add_paragraph()
p4_b2.space_after = Pt(10)
p4_b2.add_run("• Specialty Join-Table: ").font.bold = True
p4_b2.add_run("Utilizes 'tblstylist_services' to resolve the Many-to-Many relationship between stylists and salon services.\n\n")
p4_b3 = tf4.add_paragraph()
p4_b3.space_after = Pt(10)
p4_b3.add_run("• Double Approval Columns: ").font.bold = True
p4_b3.add_run("Links admin 'Status' and stylist 'StylistStatus' in 'tblappointment' for coordinated workflows.")
for para in [p4_b1, p4_b2, p4_b3]:
    for run in para.runs:
        run.font.size = Pt(12)
        run.font.color.rgb = COLOR_TEXT_MAIN

# Right Box (Structured list of tables)
create_card_shape(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
tf4_c = slide4.shapes[-1].text_frame
tf4_c.word_wrap = True
tf4_c.margin_left = Inches(0.3)
tf4_c.margin_top = Inches(0.3)
p4_ct = tf4_c.paragraphs[0]
p4_ct.add_run("Core Tables & Purposes\n\n").font.size = Pt(16)
p4_ct.runs[0].font.bold = True
p4_ct.runs[0].font.color.rgb = COLOR_SECONDARY

def add_table_item(tf, tbl_name, tbl_purpose):
    p = tf.add_paragraph()
    p.space_after = Pt(8)
    run1 = p.add_run(f"•  {tbl_name}: ")
    run1.font.bold = True
    run1.font.name = "Consolas"
    run1.font.size = Pt(11)
    run1.font.color.rgb = COLOR_PRIMARY
    
    run2 = p.add_run(tbl_purpose)
    run2.font.name = FONT_BODY
    run2.font.size = Pt(11)
    run2.font.color.rgb = COLOR_TEXT_MAIN

add_table_item(tf4_c, "tbladmin", "System administrator access credentials.")
add_table_item(tf4_c, "tblcustomers", "Client registration profiles & logins.")
add_table_item(tf4_c, "tblstylists", "Stylist profile information and credentials.")
add_table_item(tf4_c, "tblservices", "List of haircut, facial, and spa services and costs.")
add_table_item(tf4_c, "tblstylist_services", "Expertise links (Many-to-Many resolving table).")
add_table_item(tf4_c, "tblappointment", "Appointment queues with status flags.")
add_table_item(tf4_c, "tblinvoice", "Customer invoice logs and receipts.")

# ==============================================================================
# SLIDE 5: Client Portal Features
# ==============================================================================
slide5 = prs.slides.add_slide(slide_layout)
add_header(slide5, "Client (Customer) Portal Features")

# Left Column (Features list)
left_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf5 = left_box5.text_frame
tf5.word_wrap = True
p5 = tf5.paragraphs[0]
p5.add_run("Empowering Self-Service Bookings\n\n").font.size = Pt(18)
p5.runs[0].font.bold = True
p5.runs[0].font.color.rgb = COLOR_PRIMARY

add_bullet(tf5, "Self-Registration:", "Allows clients to input profiles, select credentials, and log in securely.")
add_bullet(tf5, "Dynamic Bookings:", "Filters stylist selection options to show only specialists mapped to the chosen service.")
add_bullet(tf5, "Real-Time Tracking:", "Check appointment status logs showing pending/approved stages from both admin and stylist.")
add_bullet(tf5, "Invoice History Ledger:", "Look up and review invoice details from past salon appointments.")

# Right Column (Stylist Specialty matching logic code box)
create_card_shape(slide5, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
tf5_c = slide5.shapes[-1].text_frame
tf5_c.word_wrap = True
tf5_c.margin_left = Inches(0.3)
tf5_c.margin_right = Inches(0.3)
tf5_c.margin_top = Inches(0.3)
p5_ct = tf5_c.paragraphs[0]
p5_ct.add_run("Technical Implementation\n\n").font.size = Pt(16)
p5_ct.runs[0].font.bold = True
p5_ct.runs[0].font.color.rgb = COLOR_SECONDARY

p5_cd = tf5_c.add_paragraph()
p5_cd.add_run("When a customer chooses a service, the system uses an AJAX call to load only matching specialists from the database:\n\n")
p5_cd.font.size = Pt(11)
p5_cd.font.color.rgb = COLOR_TEXT_MUTED

# Code block mock card
p5_code = tf5_c.add_paragraph()
p5_code.add_run(
    "SELECT s.ID, s.StylistName \n"
    "FROM tblstylists s \n"
    "JOIN tblstylist_services ss ON s.ID = ss.StylistId \n"
    "WHERE ss.ServiceId = :selected_service_id"
)
p5_code.font.name = "Consolas"
p5_code.font.size = Pt(11)
p5_code.font.color.rgb = COLOR_PRIMARY

# ==============================================================================
# SLIDE 6: Stylist Panel & Workflows
# ==============================================================================
slide6 = prs.slides.add_slide(slide_layout)
add_header(slide6, "Stylist Panel & Booking Approvals")

# Left Column (Double-Verification card)
create_card_shape(slide6, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
tf6_c = slide6.shapes[-1].text_frame
tf6_c.word_wrap = True
tf6_c.margin_left = Inches(0.3)
tf6_c.margin_top = Inches(0.3)
p6_ct = tf6_c.paragraphs[0]
p6_ct.add_run("Double-Verification Workflow\n\n").font.size = Pt(16)
p6_ct.runs[0].font.bold = True
p6_ct.runs[0].font.color.rgb = COLOR_PRIMARY

p6_cd = tf6_c.add_paragraph()
p6_cd.add_run(
    "To coordinate schedules, appointments utilize double-verification flags:\n\n"
    "1. Admin Approval Flag (Status)\n"
    "   • 1 = Accepted, 2 = Rejected, empty = Pending\n\n"
    "2. Stylist Approval Flag (StylistStatus)\n"
    "   • 1 = Accepted, 2 = Rejected, empty = Pending\n\n"
    "• Overall Status: Confirmed ONLY if BOTH are accepted. If either rejects, the appointment is marked Rejected."
)
for run in p6_cd.runs:
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_TEXT_MAIN

# Right Column (Text info)
right_box6 = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
tf6_r = right_box6.text_frame
tf6_r.word_wrap = True
p6_rt = tf6_r.paragraphs[0]
p6_rt.add_run("Stylist Operations\n\n").font.size = Pt(18)
p6_rt.runs[0].font.bold = True
p6_rt.runs[0].font.color.rgb = COLOR_SECONDARY

add_bullet(tf6_r, "My Schedule Dashboard:", "View client bookings containing requests matching their specialty.")
add_bullet(tf6_r, "Schedule Response Hub:", "Update booking status parameters and add comments (StylistRemark).")
add_bullet(tf6_r, "Self-Profile Controls:", "Update contact details and edit specialty details directly.")

# ==============================================================================
# SLIDE 7: Admin Panel Operations
# ==============================================================================
slide7 = prs.slides.add_slide(slide_layout)
add_header(slide7, "Admin Dashboard & Control Console")

# Subheader
sub_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.6))
sub_tf7 = sub_box7.text_frame
sub_tf7.word_wrap = True
sub_p7 = sub_tf7.paragraphs[0]
sub_r7 = sub_p7.add_run()
sub_r7.text = "Central management center tracking catalogs, staff expertise, billing, and reports."
sub_r7.font.name = FONT_BODY
sub_r7.font.size = Pt(14)
sub_r7.font.color.rgb = COLOR_TEXT_MUTED

# 3 Columns for features
col_w = Inches(3.7)
col_h = Inches(4.5)
gaps = Inches(0.3)
start_x = Inches(0.8)
start_y = Inches(2.0)

# Card 1: Dashboard Analytics
create_card_shape(slide7, start_x, start_y, col_w, col_h)
tf7_c1 = slide7.shapes[-1].text_frame
tf7_c1.word_wrap = True
tf7_c1.margin_left = Inches(0.25)
tf7_c1.margin_top = Inches(0.3)
p7_card1 = tf7_c1.paragraphs[0]
p7_card1.add_run("📊  Analytics Dashboard\n\n").font.size = Pt(16)
p7_card1.runs[0].font.bold = True
p7_card1.runs[0].font.color.rgb = COLOR_PRIMARY
p7_card1.add_run("• Real-Time Statistics: ").font.bold = True
p7_card1.add_run("Aggregates totals of active customers, pending slots, and services catalog.\n\n")
p7_card1.add_run("• Sales Tracker: ").font.bold = True
p7_card1.add_run("Calculates revenue analytics (Today, Yesterday, Last 7 Days, and Total Sales).")
for run in p7_card1.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# Card 2: Billing & Invoicing
create_card_shape(slide7, start_x + col_w + gaps, start_y, col_w, col_h)
tf7_c2 = slide7.shapes[-1].text_frame
tf7_c2.word_wrap = True
tf7_c2.margin_left = Inches(0.25)
tf7_c2.margin_top = Inches(0.3)
p7_card2 = tf7_c2.paragraphs[0]
p7_card2.add_run("🧾  Billing & Invoices\n\n").font.size = Pt(16)
p7_card2.runs[0].font.bold = True
p7_card2.runs[0].font.color.rgb = COLOR_PRIMARY
p7_card2.add_run("• Automated Billing: ").font.bold = True
p7_card2.add_run("Select a customer and check off multiple rendered services.\n\n")
p7_card2.add_run("• Print-Ready Invoices: ").font.bold = True
p7_card2.add_run("Instantly structures clean invoices equipped with a Javascript quick-print button.")
for run in p7_card2.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# Card 3: Reports & Settings
create_card_shape(slide7, start_x + (col_w + gaps)*2, start_y, col_w, col_h)
tf7_c3 = slide7.shapes[-1].text_frame
tf7_c3.word_wrap = True
tf7_c3.margin_left = Inches(0.25)
tf7_c3.margin_top = Inches(0.3)
p7_card3 = tf7_c3.paragraphs[0]
p7_card3.add_run("📈  Reports & CMS\n\n").font.size = Pt(16)
p7_card3.runs[0].font.bold = True
p7_card3.runs[0].font.color.rgb = COLOR_PRIMARY
p7_card3.add_run("• Interval Reports: ").font.bold = True
p7_card3.add_run("Filters booking records and sales stats by customized date ranges.\n\n")
p7_card3.add_run("• CMS Control: ").font.bold = True
p7_card3.add_run("Update contact info, location addresses, and operational hours dynamically.")
for run in p7_card3.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# ==============================================================================
# SLIDE 8: Interface Outputs & Showroom
# ==============================================================================
slide8 = prs.slides.add_slide(slide_layout)
add_header(slide8, "System Interfaces & Execution Outputs")

# Grid of 4 screen descriptions with metadata
card_w = Inches(5.6)
card_h = Inches(2.2)
start_x = Inches(0.8)
start_y = Inches(1.8)
gap_x = Inches(0.5)
gap_y = Inches(0.4)

def add_showroom_card(slide, x, y, title, desc, img_ref):
    create_card_shape(slide, x, y, card_w, card_h)
    tf = slide.shapes[-1].text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    
    p = tf.paragraphs[0]
    p.add_run(f"🖼️  {title}\n\n").font.bold = True
    p.runs[0].font.color.rgb = COLOR_PRIMARY
    p.runs[0].font.size = Pt(14)
    
    p.add_run(desc).font.size = Pt(11)
    p.runs[1].font.color.rgb = COLOR_TEXT_MAIN
    
    p.add_run(f"\n[Image: WORKING IMAGES/{img_ref}]").font.size = Pt(9.5)
    p.runs[2].font.color.rgb = COLOR_TEXT_MUTED

add_showroom_card(slide8, start_x, start_y, "Admin Statistics Panel", "Displays total customers, total appointments count, accepted/rejected totals, and service statistics.", "Screenshot 5.png")
add_showroom_card(slide8, start_x + card_w + gap_x, start_y, "Digital Invoicing Interface", "Lists customer details alongside costs. Incorporates an automated print mechanism.", "Screenshot 8.png")
add_showroom_card(slide8, start_x, start_y + card_h + gap_y, "Client Booking Form", "Client-side appointment registration page. Includes service selection dropdowns.", "Screenshot 10.png")
add_showroom_card(slide8, start_x + card_w + gap_x, start_y + card_h + gap_y, "Pricing & Services Catalog", "Visitor catalog showing descriptions and prices (e.g. Normal Pedicure: $400, Style Haircut: $550).", "Screenshot 11.png")

# ==============================================================================
# SLIDE 9: System Workflows
# ==============================================================================
slide9 = prs.slides.add_slide(slide_layout)
add_header(slide9, "Core System Workflows")

# Left: Booking Steps
left_box9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf9 = left_box9.text_frame
tf9.word_wrap = True
p9 = tf9.paragraphs[0]
p9.add_run("Step-by-Step Scheduling Pipeline\n\n").font.size = Pt(18)
p9.runs[0].font.bold = True
p9.runs[0].font.color.rgb = COLOR_PRIMARY

add_bullet(tf9, "1. Catalog Selection:", "Visitor reviews services pricing list on the public portal.")
add_bullet(tf9, "2. Request Placement:", "Client logs in, chooses date/time slot, and picks an expert stylist.")
add_bullet(tf9, "3. Double Validation:", "Stylist accepts assigned slot; Admin approves overall scheduling availability.")
add_bullet(tf9, "4. Service Completion:", "Staff performs requested haircuts, facials, or spa treatments.")

# Right: Billing Box
create_card_shape(slide9, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
tf9_c = slide9.shapes[-1].text_frame
tf9_c.word_wrap = True
tf9_c.margin_left = Inches(0.3)
tf9_c.margin_top = Inches(0.3)
p9_ct = tf9_c.paragraphs[0]
p9_ct.add_run("Automated Invoicing Loop\n\n").font.size = Pt(16)
p9_ct.runs[0].font.bold = True
p9_ct.runs[0].font.color.rgb = COLOR_SECONDARY

p9_cd = tf9_c.add_paragraph()
p9_cd.add_run(
    "Following service delivery, the administrator closes the booking loop:\n\n"
    "• Step A: Select the customer profile.\n\n"
    "• Step B: Tick off services rendered from the checklist.\n\n"
    "• Step C: Database logs items in 'tblinvoice' under a common 'BillingId'.\n\n"
    "• Step D: System outputs a printable invoice receipt for the customer."
)
for run in p9_cd.runs:
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_TEXT_MAIN

# ==============================================================================
# SLIDE 10: Conclusion & Future Scope
# ==============================================================================
slide10 = prs.slides.add_slide(slide_layout)
add_header(slide10, "Conclusion & Future Scope")

# Title Centered
title_box10 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(1.5))
tf10 = title_box10.text_frame
tf10.word_wrap = True
p10 = tf10.paragraphs[0]
p10.alignment = PP_ALIGN.CENTER
run10 = p10.add_run("Automating Saloon Operations Successfully")
run10.font.name = FONT_TITLE
run10.font.size = Pt(26)
run10.font.bold = True
run10.font.color.rgb = COLOR_TEXT_MAIN

p10_sub = tf10.add_paragraph()
p10_sub.alignment = PP_ALIGN.CENTER
run10_sub = p10_sub.add_run("MSMS establishes an organized, paperless, and synchronized platform linking customers, stylists, and administrators.")
run10_sub.font.name = FONT_BODY
run10_sub.font.size = Pt(13)
run10_sub.font.color.rgb = COLOR_TEXT_MUTED

# Two bottom blocks
col_w10 = Inches(5.4)
col_h10 = Inches(3.0)
start_x10 = Inches(0.8)
start_y10 = Inches(3.2)

# Card Left: Project Outcomes
create_card_shape(slide10, start_x10, start_y10, col_w10, col_h10)
tf10_cl = slide10.shapes[-1].text_frame
tf10_cl.word_wrap = True
tf10_cl.margin_left = Inches(0.25)
tf10_cl.margin_top = Inches(0.25)
p10_l = tf10_cl.paragraphs[0]
p10_l.add_run("✓  Key System Deliverables\n\n").font.bold = True
p10_l.runs[0].font.color.rgb = COLOR_PRIMARY
p10_l.runs[0].font.size = Pt(14)
p10_l.add_run(
    "• Seamless self-scheduling reduces lobby wait times.\n"
    "• Structured dashboards coordinate stylist shifts.\n"
    "• Automated records provide clean auditing database tables."
)
for run in p10_l.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# Card Right: Future scope
create_card_shape(slide10, start_x10 + col_w10 + Inches(0.533), start_y10, col_w10, col_h10)
tf10_cr = slide10.shapes[-1].text_frame
tf10_cr.word_wrap = True
tf10_cr.margin_left = Inches(0.25)
tf10_cr.margin_top = Inches(0.25)
p10_r = tf10_cr.paragraphs[0]
p10_r.add_run("🚀  Planned Future Scope\n\n").font.bold = True
p10_r.runs[0].font.color.rgb = COLOR_PRIMARY
p10_r.runs[0].font.size = Pt(14)
p10_r.add_run(
    "• Payment Integration: Stripe/UPI for online deposits.\n"
    "• Notifications: Twilio API automated SMS reminders.\n"
    "• Staff Rating: Dynamic customer reviews for stylist profiles."
)
for run in p10_r.runs[1:]:
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_TEXT_MAIN

# Save the presentation
output_name = "PROJECT_PRESENTATION.pptx"
try:
    prs.save(output_name)
    print(f"\n[SUCCESS] PowerPoint presentation generated successfully: '{output_name}'")
except Exception as e:
    print(f"\n[ERROR] Failed to save presentation: {e}")
    sys.exit(1)
